#!/usr/bin/env python3
"""Generate greppable text mirrors of an agent's materials, so indexing agents can read the
cheap mirror first and reserve expensive visual PDF reads for Hebrew/scanned pages.

Mirrors <agent>/materials/**/*.{pdf,pptx} to <agent>/materials/text/**/*.txt (same layout).
PDFs via PyMuPDF (fitz); PPTX via python-pptx. RTL languages come out garbled — that's the
signal to read those pages visually instead.

Usage: python3 replication/make_text_mirrors.py <agent>/materials
Deps:  pip install pymupdf python-pptx
"""
import os, sys, glob


def mirror_path(base, src):
    out = os.path.join(base, "text", os.path.relpath(src, base))
    os.makedirs(os.path.dirname(out), exist_ok=True)
    return os.path.splitext(out)[0] + ".txt"


def main(base):
    npdf = npptx = 0
    try:
        import fitz  # PyMuPDF
        for f in glob.glob(f"{base}/**/*.pdf", recursive=True):
            if "/text/" in f:
                continue
            try:
                doc = fitz.open(f)
                txt = "\n".join(f"\n===== page {i+1} =====\n" + p.get_text()
                                for i, p in enumerate(doc))
                open(mirror_path(base, f), "w").write(txt)
                npdf += 1
            except Exception as e:
                print("PDF ERR", f, e, file=sys.stderr)
    except ImportError:
        print("PyMuPDF (fitz) not installed — skipping PDFs. pip install pymupdf", file=sys.stderr)

    try:
        from pptx import Presentation
        for f in glob.glob(f"{base}/**/*.pptx", recursive=True):
            if "/text/" in f:
                continue
            try:
                prs = Presentation(f)
                parts = []
                for i, slide in enumerate(list(prs.slides)):
                    st = [sh.text_frame.text.strip() for sh in slide.shapes
                          if sh.has_text_frame and sh.text_frame.text.strip()]
                    parts.append(f"\n===== slide {i+1} =====\n" + "\n".join(st))
                open(mirror_path(base, f), "w").write("\n".join(parts))
                npptx += 1
            except Exception as e:
                print("PPTX ERR", f, e, file=sys.stderr)
    except ImportError:
        print("python-pptx not installed — skipping PPTX. pip install python-pptx", file=sys.stderr)

    print(f"text mirrors written: {npdf} pdf, {npptx} pptx -> {base}/text/")


if __name__ == "__main__":
    if len(sys.argv) != 2:
        sys.exit("usage: make_text_mirrors.py <agent>/materials")
    main(sys.argv[1].rstrip("/"))
